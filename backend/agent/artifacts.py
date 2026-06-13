"""Binary artifact generation (docx / pptx / pdf / xlsx).

Files are written under static/files/ and served by the same backend; the agent
returns an ABSOLUTE artifact_url built from PUBLIC_BASE_URL. Inline HTML/markdown
decks are NOT handled here — those go straight into the answer text.

The model supplies the data as a structured spec (it has already gathered the
facts via tools); these functions only render it.
"""

from __future__ import annotations

import hashlib
import os
import re
from pathlib import Path
from typing import Any

_FILES = Path(__file__).resolve().parent.parent / "static" / "files"
_FILES.mkdir(parents=True, exist_ok=True)

BINARY_KINDS = ("docx", "pptx", "pdf", "xlsx")


def detect_binary_kind(question: str) -> str | None:
    """Return the explicitly-requested binary format, or None (→ inline HTML/md)."""
    q = question.lower()
    # Match whole words / extensions to avoid false hits.
    for kind in BINARY_KINDS:
        if re.search(rf"\b{kind}\b|\.{kind}\b", q):
            return kind
    if "powerpoint" in q or "slide deck file" in q:
        return "pptx"
    if "word document" in q:
        return "docx"
    if "excel" in q or "spreadsheet" in q:
        return "xlsx"
    return None


def _slug(title: str, kind: str, payload: str) -> str:
    base = re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-")[:50] or "artifact"
    h = hashlib.sha1(payload.encode("utf-8")).hexdigest()[:8]
    return f"{base}-{h}.{kind}"


def _url(filename: str) -> str:
    base = os.environ.get("PUBLIC_BASE_URL", "http://localhost:8000").rstrip("/")
    return f"{base}/files/{filename}"


# --- renderers ---------------------------------------------------------------

def _docx(path: Path, title: str, sections: list[dict[str, Any]]) -> None:
    from docx import Document
    doc = Document()
    doc.add_heading(title, level=0)
    for s in sections:
        if s.get("heading"):
            doc.add_heading(str(s["heading"]), level=1)
        body = s.get("body", "")
        for para in str(body).split("\n"):
            if para.strip():
                doc.add_paragraph(para.strip())
    doc.save(path)


def _pptx(path: Path, title: str, slides: list[dict[str, Any]]) -> None:
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    title_layout, bullet_layout = prs.slide_layouts[0], prs.slide_layouts[1]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = title
    for slide in slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = str(slide.get("title", ""))
        tf = s.placeholders[1].text_frame
        tf.clear()
        bullets = slide.get("bullets") or ([slide["body"]] if slide.get("body") else [])
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(b)
            p.font.size = Pt(18)
    prs.save(path)


def _pdf(path: Path, title: str, sections: list[dict[str, Any]]) -> None:
    from fpdf import FPDF
    pdf = FPDF()
    pdf.set_margins(15, 15, 15)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    w = pdf.epw  # effective page width (page minus L/R margins) — avoids width errors
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(w, 10, _latin1(title))
    for s in sections:
        if s.get("heading"):
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(w, 8, _latin1(str(s["heading"])))
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(w, 6, _latin1(str(s.get("body", ""))))
    pdf.output(str(path))


def _xlsx(path: Path, title: str, table: dict[str, Any], sheet_name: str | None) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = (sheet_name or title or "Sheet1")[:31]
    columns = table.get("columns") or []
    if columns:
        ws.append([str(c) for c in columns])
    for row in table.get("rows") or []:
        ws.append(list(row))
    wb.save(path)


def _latin1(text: str) -> str:
    """fpdf2 core fonts are latin-1; drop unrepresentable chars safely."""
    return text.encode("latin-1", "replace").decode("latin-1")


def create(kind: str, title: str, *, sections: list[dict] | None = None,
           slides: list[dict] | None = None, table: dict | None = None,
           sheet_name: str | None = None) -> str:
    """Render an artifact and return its absolute URL."""
    kind = kind.lower()
    payload = repr((title, sections, slides, table))
    filename = _slug(title or "artifact", kind, payload)
    path = _FILES / filename
    if kind == "docx":
        _docx(path, title, sections or [])
    elif kind == "pptx":
        _pptx(path, title, slides or [])
    elif kind == "pdf":
        _pdf(path, title, sections or [])
    elif kind == "xlsx":
        _xlsx(path, title, table or {}, sheet_name)
    else:
        raise ValueError(f"unsupported artifact kind: {kind}")
    return _url(filename)
